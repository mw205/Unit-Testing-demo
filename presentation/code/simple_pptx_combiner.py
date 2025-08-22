"""
Simple and reliable PowerPoint combiner
Install required library: pip install python-pptx
"""

from pptx import Presentation
import os

def combine_pptx_files(input_files, output_file):
    """
    Combines multiple PowerPoint files into one
    
    Args:
        input_files: List of input .pptx file paths
        output_file: Output file path for combined presentation
    """
    
    # Check if input files exist
    valid_files = []
    for file in input_files:
        if os.path.exists(file):
            valid_files.append(file)
            print(f"✓ Found: {file}")
        else:
            print(f"✗ Missing: {file}")
    
    if not valid_files:
        print("❌ No valid input files found!")
        return False
    
    print(f"\n📊 Combining {len(valid_files)} presentations...")
    
    # Use the first file as the base
    combined_presentation = Presentation(valid_files[0])
    total_slides = len(combined_presentation.slides)
    print(f"Base presentation: {valid_files[0]} ({total_slides} slides)")
    
    # Add slides from other presentations
    for file_path in valid_files[1:]:
        try:
            source_prs = Presentation(file_path)
            slides_added = 0
            
            for source_slide in source_prs.slides:
                # Get the layout (use blank layout as fallback)
                try:
                    slide_layout = source_slide.slide_layout
                except:
                    slide_layout = combined_presentation.slide_layouts[6]  # Blank layout
                
                # Add new slide to combined presentation
                new_slide = combined_presentation.slides.add_slide(slide_layout)
                
                # Copy all shapes from source slide
                for shape in source_slide.shapes:
                    # Use the shape's XML element to preserve formatting
                    el = shape.element
                    new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
                
                slides_added += 1
            
            total_slides += slides_added
            print(f"Added from {file_path}: {slides_added} slides")
            
        except Exception as e:
            print(f"⚠️  Error with {file_path}: {str(e)}")
            continue
    
    # Save the combined presentation
    try:
        combined_presentation.save(output_file)
        print(f"\n✅ Success! Combined presentation saved as: {output_file}")
        print(f"📋 Total slides: {total_slides}")
        return True
        
    except Exception as e:
        print(f"❌ Error saving file: {str(e)}")
        return False

# Main execution
if __name__ == "__main__":
    # Your three files
    input_files = [
        "Unit_Testing_Python4.pptx",
        "Unit_Testing_Python5.pptx", 
        "unittest_presentation.pptx", 
        "Unit_Testing_Python6.pptx"
    ]
    
    output_file = "Combined_Unit_Testing_Python.pptx"
    
    print("🔗 PowerPoint File Combiner")
    print("=" * 40)
    
    # Run the combination
    success = combine_pptx_files(input_files, output_file)
    
    if success:
        print(f"\n🎉 All done! Your combined presentation is ready: {output_file}")
    else:
        print("\n💥 Something went wrong. Please check the error messages above.")

# Quick usage function for interactive use
def quick_combine():
    """Quick function to combine the three specific files"""
    files = ["Unit_Testing_Python4.pptx", "Unit_Testing_Python5.pptx", "Unit_Testing_Python4.pptx", "unittest_presentation.pptx"]
    return combine_pptx_files(files, "Combined_Unit_Testing_Python.pptx")
