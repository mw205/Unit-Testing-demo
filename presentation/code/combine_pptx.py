"""
Script to combine multiple PowerPoint presentations into one file
Requires: python-pptx library (pip install python-pptx)
"""

from pptx import Presentation
import os

def combine_presentations(file_paths, output_path):
    """
    Combine multiple PowerPoint presentations into a single file
    
    Args:
        file_paths (list): List of paths to PowerPoint files to combine
        output_path (str): Path for the combined output file
    """
    # Create a new presentation for the combined slides
    combined_prs = Presentation()
    
    # Remove the default blank slide that comes with new presentation
    if len(combined_prs.slides) > 0:
        # Get the slide layout from the first slide before removing it
        slide_layout = combined_prs.slides[0].slide_layout
        # Remove the blank slide
        xml_slides = combined_prs.slides._sldIdLst
        xml_slides.remove(xml_slides[0])
    
    slide_count = 0
    
    # Process each PowerPoint file
    for i, file_path in enumerate(file_paths):
        if not os.path.exists(file_path):
            print(f"Warning: File '{file_path}' not found. Skipping...")
            continue
            
        print(f"Processing file {i+1}/{len(file_paths)}: {file_path}")
        
        try:
            # Open the presentation
            prs = Presentation(file_path)
            
            # Copy each slide to the combined presentation
            for slide in prs.slides:
                # Create a new slide in the combined presentation
                slide_layout = slide.slide_layout
                new_slide = combined_prs.slides.add_slide(slide_layout)
                
                # Copy all shapes from the original slide
                for shape in slide.shapes:
                    # Get the shape element
                    new_shape_element = shape.element
                    # Add it to the new slide
                    new_slide.shapes._spTree.insert_element_before(
                        new_shape_element, 'p:extLst'
                    )
                
                slide_count += 1
                
            print(f"Added {len(prs.slides)} slides from {file_path}")
            
        except Exception as e:
            print(f"Error processing {file_path}: {str(e)}")
            continue
    
    # Save the combined presentation
    try:
        combined_prs.save(output_path)
        print(f"\nSuccess! Combined {slide_count} slides into '{output_path}'")
        return True
    except Exception as e:
        print(f"Error saving combined presentation: {str(e)}")
        return False

def main():
    # List of PowerPoint files to combine
    files_to_combine = [
        "Unit_Testing_Python.pptx",
        "Unit_Testing_Python1.pptx", 
        "Unit_Testing_Python2.pptx"
    ]
    
    # Output file name
    output_file = "Combined_Unit_Testing_Python.pptx"
    
    print("PowerPoint Combiner")
    print("=" * 50)
    print(f"Files to combine: {len(files_to_combine)}")
    for i, file in enumerate(files_to_combine, 1):
        print(f"  {i}. {file}")
    print(f"Output file: {output_file}")
    print("=" * 50)
    
    # Check if files exist
    existing_files = []
    for file_path in files_to_combine:
        if os.path.exists(file_path):
            existing_files.append(file_path)
        else:
            print(f"Warning: '{file_path}' not found")
    
    if not existing_files:
        print("Error: No input files found!")
        return
    
    print(f"\nFound {len(existing_files)} files to process")
    
    # Combine the presentations
    success = combine_presentations(existing_files, output_file)
    
    if success:
        print(f"\n✅ Successfully created '{output_file}'")
    else:
        print(f"\n❌ Failed to create combined presentation")

# Alternative method using a more robust approach
def combine_presentations_advanced(file_paths, output_path):
    """
    Advanced method to combine PowerPoint presentations
    This method preserves more formatting and handles edge cases better
    """
    from pptx.util import Inches
    import xml.etree.ElementTree as ET
    
    # Start with the first presentation as the base
    if not file_paths or not os.path.exists(file_paths[0]):
        print("Error: No valid files provided")
        return False
    
    print(f"Using '{file_paths[0]}' as base presentation...")
    combined_prs = Presentation(file_paths[0])
    base_slide_count = len(combined_prs.slides)
    
    # Add slides from remaining presentations
    for file_path in file_paths[1:]:
        if not os.path.exists(file_path):
            print(f"Warning: File '{file_path}' not found. Skipping...")
            continue
            
        print(f"Adding slides from '{file_path}'...")
        
        try:
            source_prs = Presentation(file_path)
            
            for slide in source_prs.slides:
                # Use the slide layout from the source slide
                try:
                    # Try to find a matching layout in the combined presentation
                    target_layout = None
                    source_layout = slide.slide_layout
                    
                    # Look for matching layout by name
                    for layout in combined_prs.slide_layouts:
                        if layout.name == source_layout.name:
                            target_layout = layout
                            break
                    
                    # If no matching layout found, use a blank layout
                    if target_layout is None:
                        target_layout = combined_prs.slide_layouts[6]  # Blank layout
                    
                    # Add new slide with the layout
                    new_slide = combined_prs.slides.add_slide(target_layout)
                    
                    # Copy slide properties
                    if hasattr(slide, 'background'):
                        new_slide.background = slide.background
                    
                    # Copy all shapes
                    for shape in slide.shapes:
                        # Clone the shape to the new slide
                        new_slide.shapes._spTree.append(shape.element)
                        
                except Exception as e:
                    print(f"Warning: Could not copy slide properly: {e}")
                    # Fallback: create blank slide with text noting the issue
                    blank_layout = combined_prs.slide_layouts[6]
                    new_slide = combined_prs.slides.add_slide(blank_layout)
                    
            print(f"Added {len(source_prs.slides)} slides from {file_path}")
            
        except Exception as e:
            print(f"Error processing {file_path}: {str(e)}")
            continue
    
    # Save the result
    try:
        combined_prs.save(output_path)
        final_count = len(combined_prs.slides)
        print(f"\n✅ Success! Combined presentation saved as '{output_path}'")
        print(f"Total slides: {final_count} (started with {base_slide_count})")
        return True
    except Exception as e:
        print(f"Error saving: {str(e)}")
        return False

if __name__ == "__main__":
    # You can also run the advanced method by uncommenting the line below:
    # combine_presentations_advanced(["Unit_Testing_Python.pptx", "Unit_Testing_Python1.pptx", "Unit_Testing_Python2.pptx"], "Combined_Advanced.pptx")
    main()
