from pptx import Presentation
from pptx.util import Inches

# Create a PowerPoint presentation
prs = Presentation()

# Define slides content
slides_content = [
    ("Mastering Unit Testing in Python with unittest", "Write Robust and Reliable Code\nYour Name/Affiliation\nDate"),
    ("What is Unit Testing?", "- Testing individual units of code (functions, methods) in isolation.\n- Focuses on verifying the smallest testable parts of an application.\n- Like testing individual gears in a machine, not the whole machine at once."),
    ("Why Unit Testing?", "- Reduces Defects: Find bugs early, saving time and money.\n- Improves Code Quality: Forces better design and modularity.\n- Facilitates Refactoring: Provides a safety net for code changes.\n- Living Documentation: Tests show how code should work.\n- Increased Confidence: Make changes without fear of breaking things."),
    ("The Testing Pyramid", "- Unit tests should form the foundation of your testing strategy.\n- Integration tests check how units work together.\n- E2E tests validate the entire application flow."),
    ("FIRST Principles of Good Unit Tests", "- Fast: Run quickly.\n- Isolated: Independent of each other.\n- Repeatable: Same result every time.\n- Self-validating: Clear pass/fail.\n- Timely: Written before/alongside code (TDD)."),
    ("Introducing unittest", "- Python's built-in testing framework.\n- Inspired by JUnit (Java).\n- Provides tools for:\n  - Creating test cases (unittest.TestCase).\n  - Making assertions (e.g., assertEqual).\n  - Running tests (test runners)."),
    ("Your First Unit Test", "Code Example:\n\ndef add(x, y):\n    return x + y\n\nimport unittest\nfrom my_module import add\n\nclass TestAdd(unittest.TestCase):\n    def test_add(self):\n        self.assertEqual(add(2, 3), 5)\n\nif __name__ == '__main__':\n    unittest.main()"),
    ("Running Tests", "- From the Command Line:\n  - python -m unittest test_my_module.py\n  - python -m unittest discover\n  - python -m unittest -v test_my_module.py (verbose)\n- From within File:\n  - if __name__ == '__main__': unittest.main()"),
    ("Assertion Methods", "- self.assertEqual(a, b): Checks a == b\n- self.assertTrue(x): Checks bool(x) is True\n- self.assertIsNone(x): Checks x is None\n- self.assertIn(a, b): Checks a in b\n- self.assertRaises(Exception, callable, *args, **kwargs): Checks for exceptions"),
    ("Mocking - Isolating Dependencies", "- Problem: Testing code with external dependencies (network, database, etc.).\n- Solution: Replace dependencies with \"test doubles\" (mocks, stubs).\n- Benefits:\n  - Faster tests.\n  - More reliable tests (no external dependencies).\n  - Control over dependency behavior."),
    ("Test-Driven Development (TDD)", "- Red-Green-Refactor:\n  - Red: Write a failing test first.\n  - Green: Write the minimum code to pass the test.\n  - Refactor: Clean up and improve the code.\n- Benefits:\n  - 100% test coverage (in theory).\n  - Better code design.\n  - Immediate feedback."),
    ("Key Takeaways", "- Unit testing is crucial for writing robust and reliable code.\n- unittest is Python's built-in, powerful testing framework.\n- Use assertions to verify expected behavior.\n- Organize your tests for maintainability.\n- Mock dependencies to isolate your code.\n- Practice TDD for better design and coverage."),
    ("Questions?", "Any questions? Feel free to ask!"),
    ("Thank You!", "Links to Resources:\n- unittest documentation: https://docs.python.org/3/library/unittest.html\n- unittest.mock documentation: https://docs.python.org/3/library/unittest.mock.html\n- (Optional) Your GitHub repository with code examples")
]

# Add slides to presentation
for title, content in slides_content:
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = content

# Save the presentation file
pptx_path = "/mnt/data/unittest_presentation.pptx"
prs.save(pptx_path)

# Provide download link
pptx_path

# Take multiple files as input and combine all the information into a single Markdown file. Organize the content logically, ensuring sections flow coherently. Refactor any redundant, inconsistent, or unclear parts to improve readability and structure. Use appropriate Markdown formatting (headings, lists, tables, and code blocks where necessary). Ensure the final document is clean, professional, and easy to navigate.
